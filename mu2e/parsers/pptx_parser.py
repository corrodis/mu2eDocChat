"""
PPTX (PowerPoint) parser
"""

import io
from PIL import Image
from pptx import Presentation
from tqdm import tqdm
from .base_parser import BaseParser

class PPTXParser(BaseParser):
    """Parser for PPTX files"""
    
    def get_text(self, rescale_image_max_dim=500):
        """Extract text and images from PPTX document"""
        extracted_text = ""
        images = []
        image_cnt = 0
        
        prs = Presentation(self.doc)
        
        for slide_num, slide in enumerate(tqdm(prs.slides, desc="Processing slides")):
            slide_text = ""
            
            # Extract text from all shapes
            for shape in slide.shapes:
                # Text boxes and title/content placeholders
                if hasattr(shape, "text_frame") and shape.text_frame.text.strip():
                    shape_text = self._extract_text_with_hyperlinks(shape.text_frame)
                    if shape_text.strip():
                        slide_text += shape_text + "\n"
                elif hasattr(shape, "text") and shape.text.strip():
                    slide_text += shape.text + "\n"
                
                # Tables
                elif hasattr(shape, "table"):
                    table_text = "\n**Table:**\n"
                    for row in shape.table.rows:
                        row_text = " | ".join([cell.text.strip() for cell in row.cells])
                        table_text += f"| {row_text} |\n"
                    slide_text += table_text + "\n"
            
            # Extract speaker notes if available
            if slide.notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_text += f"\n**Speaker Notes:**\n{notes_text}\n"
            
            # Clean and format the text
            cleaned_text = self._clean_text(slide_text)
            markdown_text = self._slides_format_as_markdown(cleaned_text)
            
            if slide_text.strip():
                extracted_text += f"<slide number={slide_num+1}>{markdown_text}"
            
            # Extract images
            for shape in slide.shapes:
                if shape.shape_type == 13:  # Picture shape type
                    try:
                        image_stream = io.BytesIO(shape.image.blob)
                        img = Image.open(image_stream)
                        img = self._resize_image(img, rescale_image_max_dim)
                        img_base64 = self._image_to_base64(img, img.format)
                        
                        images.append(img_base64)
                        extracted_text += f"[Image {image_cnt}]"
                        image_cnt += 1
                    except Exception as e:
                        print(f"Error processing image in slide {slide_num+1}: {e}")
        
        return extracted_text, images
    
    def _extract_text_with_hyperlinks(self, text_frame):
        """Extract text including hyperlinks from text frame"""
        full_text = ""
        for paragraph in text_frame.paragraphs:
            para_text = ""
            for run in paragraph.runs:
                run_text = run.text
                if run.hyperlink and run.hyperlink.address:
                    para_text += f"[{run_text}]({run.hyperlink.address})"
                else:
                    para_text += run_text
            full_text += para_text + "\n"
        return full_text